"""Tests for PPTX full extraction."""

import os
import tempfile
from runeextract import extract


def test_pptx_extraction():
    """Test basic PPTX text extraction."""
    from pptx import Presentation
    from pptx.util import Inches
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        path = tmp.name
    try:
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # title slide layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "My Title"
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        tf = txBox.text_frame
        tf.text = "Hello from python-pptx."
        prs.save(path)

        result = extract(path)
        assert result.source_type == "pptx"
        assert "My Title" in result.text or "Hello from python-pptx" in result.text
    finally:
        if os.path.exists(path):
            os.unlink(path)


def test_pptx_with_table():
    """Test PPTX with table extraction."""
    from pptx import Presentation
    from pptx.util import Inches
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        path = tmp.name
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        table_shape = slide.shapes.add_table(2, 3, Inches(1), Inches(1), Inches(6), Inches(2))
        table = table_shape.table
        table.cell(0, 0).text = "A"
        table.cell(0, 1).text = "B"
        table.cell(0, 2).text = "C"
        table.cell(1, 0).text = "1"
        table.cell(1, 1).text = "2"
        table.cell(1, 2).text = "3"
        prs.save(path)

        result = extract(path)
        assert len(result.tables) > 0
    finally:
        if os.path.exists(path):
            os.unlink(path)


def test_pptx_with_metadata():
    """Test PPTX metadata extraction."""
    from pptx import Presentation
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        path = tmp.name
    try:
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.core_properties.author = "Test Author"
        prs.core_properties.title = "Test PPTX"
        prs.save(path)

        result = extract(path, metadata=True)
        assert result.metadata.get("title") == "Test PPTX"
    finally:
        if os.path.exists(path):
            os.unlink(path)
